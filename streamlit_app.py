"""
راج للذكاء الاصطناعي للوثائق - تطبيق محسن
الميزات الجديدة:
- واجهة عربية كاملة
- دعم OCR للصور
- بحث صوتي
- تحليلات متقدمة
- معالجة فورية
- تصدير محسن
- تحليل تشابه الوثائق
- اقتراحات ذكية
- تتبع التقدم
- تحليلات بصرية
- حفظ الجلسات
- تحسينات الأداء
"""
import os
import io
import time
import math
import tempfile
import traceback
import json
import base64
import hashlib
import pickle
import re
from typing import List, Dict, Tuple, Optional, Any
from datetime import datetime, timedelta
from collections import defaultdict, Counter

import streamlit as st
import pandas as pd
import numpy as np

# الوارداات المحسنة
try:
    from googletrans import Translator
    TRANSLATOR_AVAILABLE = True
except Exception:
    TRANSLATOR_AVAILABLE = False

try:
    import speech_recognition as sr
    SPEECH_RECOGNITION_AVAILABLE = True
except Exception:
    SPEECH_RECOGNITION_AVAILABLE = False

try:
    import plotly.express as px
    import plotly.graph_objects as go
    from plotly.subplots import make_subplots
    PLOTLY_AVAILABLE = True
except Exception:
    PLOTLY_AVAILABLE = False

try:
    from PIL import Image
    import pytesseract
    OCR_AVAILABLE = True
except Exception:
    OCR_AVAILABLE = False

try:
    import cv2
    OPENCV_AVAILABLE = True
except Exception:
    OPENCV_AVAILABLE = False

try:
    import networkx as nx
    NETWORKX_AVAILABLE = True
except Exception:
    NETWORKX_AVAILABLE = False

try:
    from wordcloud import WordCloud
    import arabic_reshaper
    from bidi.algorithm import get_display
    WORDCLOUD_AVAILABLE = True
except Exception:
    WORDCLOUD_AVAILABLE = False

try:
    import textstat
    TEXTSTAT_AVAILABLE = True
except Exception:
    TEXTSTAT_AVAILABLE = False

# الوارداات الأصلية
try:
    from modules.file_handler import load_file
    from modules.utils import clean_text, chunk_text, sentences_from_text, tfidf_sentence_ranking
    from modules.ai_engine import RAGEngine, get_embedding as module_get_embedding, chat_with_ai as module_chat_with_ai
    from modules.exporter import export_txt, export_docx, export_pdf
    from modules.ui_components import render_upload_box, render_chat_ui
    MODULES_AVAILABLE = True
except Exception:
    MODULES_AVAILABLE = False

try:
    from sentence_transformers import SentenceTransformer
except Exception:
    SentenceTransformer = None

try:
    import faiss
except Exception:
    faiss = None

try:
    from openai import OpenAI
except Exception:
    OpenAI = None

try:
    from PyPDF2 import PdfReader
except Exception:
    PdfReader = None

try:
    import docx
except Exception:
    docx = None

try:
    import openpyxl
except Exception:
    openpyxl = None

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
except Exception:
    Presentation = None

try:
    from docx import Document as DocxDoc
except Exception:
    DocxDoc = None

try:
    import whisper
    WHISPER_AVAILABLE = True
except Exception:
    WHISPER_AVAILABLE = False

# ---------------------- الإعدادات والأسرار ----------------------
st.set_page_config(
    page_title="راج للذكاء الاصطناعي للوثائق", 
    page_icon="🌍", 
    layout="wide", 
    initial_sidebar_state="expanded"
)

# CSS مخصص للواجهة العربية
st.markdown("""
<style>
    .rtl-text {
        direction: rtl;
        text-align: right;
        font-family: 'Arial', sans-serif;
    }
    .metric-card {
        background-color: #f0f2f6;
        padding: 1rem;
        border-radius: 0.5rem;
        border-left: 4px solid #1f77b4;
        margin: 0.5rem 0;
    }
    .success-card {
        background-color: #d4edda;
        border-color: #28a745;
    }
    .warning-card {
        background-color: #fff3cd;
        border-color: #ffc107;
    }
    .error-card {
        background-color: #f8d7da;
        border-color: #dc3545;
    }
    .sidebar .sidebar-content {
        background-image: linear-gradient(#2e7bcf,#2e7bcf);
    }
    .stProgress .st-bo {
        background-color: #1f77b4;
    }
    .chat-message {
        padding: 1rem;
        border-radius: 0.5rem;
        margin: 0.5rem 0;
    }
    .user-message {
        background-color: #e3f2fd;
        margin-left: 2rem;
    }
    .assistant-message {
        background-color: #f5f5f5;
        margin-right: 2rem;
    }
</style>
""", unsafe_allow_html=True)

st.title("🌍 راج للذكاء الاصطناعي للوثائق — منصة RAG المتطورة")

# قراءة مفتاح OpenAI
OPENAI_API_KEY = None
try:
    OPENAI_API_KEY = st.secrets.get("OPENAI_API_KEY")
except Exception:
    pass
if not OPENAI_API_KEY:
    OPENAI_API_KEY = os.environ.get("OPENAI_API_KEY")

# إنشاء عميل OpenAI
openai_client = None
if OPENAI_API_KEY and OpenAI is not None:
    try:
        openai_client = OpenAI(api_key=OPENAI_API_KEY)
    except Exception:
        openai_client = None

# إنشاء مترجم
translator = None
if TRANSLATOR_AVAILABLE:
    try:
        translator = Translator()
    except Exception:
        translator = None

# ---------------------- إدارة الحالة المحسنة ----------------------
def initialize_session_state():
    """تهيئة حالة الجلسة مع القيم الافتراضية"""
    defaults = {
        "docs": [],
        "index_built": False,
        "faiss_index": None,
        "embeddings": None,
        "id_map": [],
        "model_name": "all-MiniLM-L6-v2",
        "chat_history": [],
        "query_count": 0,
        "processing_stats": {
            "total_docs": 0,
            "total_chars": 0,
            "total_words": 0,
            "avg_query_time": 0,
            "languages_detected": set(),
            "file_types": {},
            "sessions_count": 0,
            "total_queries": 0,
            "successful_queries": 0
        },
        "user_preferences": {
            "language": "ar",
            "theme": "light",
            "chunk_size": 300,
            "chunk_overlap": 50,
            "top_k_results": 5,
            "temperature": 0.7,
            "max_tokens": 1000
        },
        "current_session": {
            "start_time": datetime.now(),
            "queries": [],
            "docs_processed": 0
        },
        "advanced_analytics": {
            "query_patterns": defaultdict(int),
            "response_ratings": [],
            "most_accessed_docs": defaultdict(int),
            "search_history": [],
            "user_feedback": []
        },
        "cache": {},
        "bookmarks": [],
        "export_history": []
    }
    
    for key, value in defaults.items():
        if key not in st.session_state:
            st.session_state[key] = value

initialize_session_state()

# ---------------------- الوظائف المساعدة المحسنة ----------------------
def app_log(msg: str):
    """تسجيل محسن مع الطابع الزمني والمستوى"""
    timestamp = datetime.now().strftime('%Y-%m-%d %H:%M:%S')
    log_entry = {
        "timestamp": timestamp,
        "message": msg,
        "level": "INFO"
    }
    st.session_state.setdefault("_logs", []).append(log_entry)
    if len(st.session_state["_logs"]) > 100:  # الحد من حجم السجل
        st.session_state["_logs"] = st.session_state["_logs"][-100:]

def generate_doc_hash(content: str) -> str:
    """توليد hash للوثيقة للتحقق من التكرار"""
    return hashlib.md5(content.encode()).hexdigest()

def detect_language(text: str) -> str:
    """كشف لغة النص محسن"""
    if not translator:
        # كشف بسيط للعربية
        arabic_chars = sum(1 for c in text[:200] if '\u0600' <= c <= '\u06FF')
        if arabic_chars > 10:
            return "ar"
        return "en"
    
    try:
        result = translator.detect(text[:200])
        st.session_state.processing_stats["languages_detected"].add(result.lang)
        return result.lang
    except Exception:
        return "unknown"

def translate_text(text: str, target_lang: str = "ar") -> str:
    """ترجمة النص محسنة مع تخزين مؤقت"""
    if not translator:
        return text
    
    # التحقق من التخزين المؤقت
    cache_key = f"translate_{hash(text)}_{target_lang}"
    if cache_key in st.session_state.cache:
        return st.session_state.cache[cache_key]
    
    try:
        result = translator.translate(text, dest=target_lang)
        translated = result.text
        st.session_state.cache[cache_key] = translated
        return translated
    except Exception as e:
        app_log(f"خطأ في الترجمة: {e}")
        return text

def extract_text_from_image(image) -> str:
    """استخراج النص من الصور باستخدام OCR محسن"""
    if not OCR_AVAILABLE:
        return "[OCR غير متاح: تثبيت pytesseract مطلوب]"
    
    try:
        # تحويل الصورة لـ PIL
        if isinstance(image, bytes):
            img = Image.open(io.BytesIO(image))
        else:
            img = Image.open(image)
        
        # تحسين الصورة للـ OCR
        if OPENCV_AVAILABLE:
            img_array = np.array(img)
            if len(img_array.shape) == 3:
                img_array = cv2.cvtColor(img_array, cv2.COLOR_RGB2GRAY)
            
            # تحسين التباين
            img_array = cv2.equalizeHist(img_array)
            img = Image.fromarray(img_array)
        
        # استخراج النص مع دعم متعدد اللغات
        text = pytesseract.image_to_string(img, lang='ara+eng+fra')
        
        # تنظيف النص المستخرج
        text = re.sub(r'\s+', ' ', text).strip()
        
        return text
    except Exception as e:
        app_log(f"خطأ في استخراج النص من الصورة: {e}")
        return f"[خطأ في استخراج النص من الصورة: {e}]"

def process_audio_input(audio_file) -> str:
    """معالجة الصوت واستخراج النص"""
    if not WHISPER_AVAILABLE and not SPEECH_RECOGNITION_AVAILABLE:
        return "[معالجة الصوت غير متاحة]"
    
    try:
        # حفظ الملف الصوتي مؤقتاً
        with tempfile.NamedTemporaryFile(delete=False, suffix=".wav") as tmp_file:
            tmp_file.write(audio_file.getvalue())
            tmp_path = tmp_file.name
        
        if WHISPER_AVAILABLE:
            # استخدام Whisper للتعرف على الصوت
            model = whisper.load_model("base")
            result = model.transcribe(tmp_path, language="ar")
            text = result["text"]
        elif SPEECH_RECOGNITION_AVAILABLE:
            # استخدام speech_recognition
            r = sr.Recognizer()
            with sr.AudioFile(tmp_path) as source:
                audio = r.record(source)
                text = r.recognize_google(audio, language='ar-SA')
        else:
            text = "[فشل في معالجة الصوت]"
        
        # تنظيف الملف المؤقت
        os.unlink(tmp_path)
        
        return text
    except Exception as e:
        app_log(f"خطأ في معالجة الصوت: {e}")
        return f"[خطأ في معالجة الصوت: {e}]"

def analyze_text_statistics(text: str) -> Dict[str, Any]:
    """تحليل إحصائيات النص المتقدمة"""
    stats = {
        "char_count": len(text),
        "word_count": len(text.split()),
        "sentence_count": len(fallback_sentences_from_text(text)),
        "paragraph_count": len([p for p in text.split('\n\n') if p.strip()]),
        "language": detect_language(text)
    }
    
    if TEXTSTAT_AVAILABLE:
        try:
            stats.update({
                "readability_score": textstat.flesch_reading_ease(text),
                "grade_level": textstat.flesch_kincaid_grade(text),
                "avg_sentence_length": textstat.avg_sentence_length(text),
                "difficult_words": textstat.difficult_words(text)
            })
        except Exception:
            pass
    
    # تحليل الكلمات الأكثر شيوعاً
    words = re.findall(r'\b\w+\b', text.lower())
    word_freq = Counter(words)
    stats["most_common_words"] = word_freq.most_common(10)
    
    return stats

def fallback_load_file_bytes(fileobj) -> Tuple[str, Dict[str, Any]]:
    """قارئ احتياطي للملفات المرفوعة مع معلومات إضافية"""
    name = getattr(fileobj, "name", "مرفوع")
    ext = os.path.splitext(name)[1].lower()
    
    try:
        raw = fileobj.read()
    except Exception as e:
        return f"[خطأ في قراءة الملف: {e}]", {"error": str(e)}
    
    # إحصائيات الملف
    file_size = len(raw)
    file_info = {
        "name": name,
        "extension": ext,
        "size_bytes": file_size,
        "size_mb": file_size / (1024 * 1024),
        "processed_at": datetime.now().isoformat()
    }
    
    st.session_state.processing_stats["file_types"][ext] = st.session_state.processing_stats["file_types"].get(ext, 0) + 1
    
    content = ""
    
    if ext == ".pdf":
        if PdfReader is None:
            content = "[قراءة PDF غير متاحة: تثبيت PyPDF2 مطلوب]"
        else:
            try:
                reader = PdfReader(io.BytesIO(raw))
                pages = []
                file_info["pages"] = len(reader.pages)
                
                for i, page in enumerate(reader.pages):
                    txt = page.extract_text() or ""
                    if txt.strip():
                        pages.append(f"--- صفحة {i+1} ---\n{txt}")
                
                content = "\n\n".join(pages)
                file_info["extracted_pages"] = len(pages)
            except Exception as e:
                content = f"[خطأ في قراءة PDF: {e}]"
                file_info["error"] = str(e)
    
    elif ext == ".docx":
        if docx is None:
            content = "[قراءة DOCX غير متاحة: تثبيت python-docx مطلوب]"
        else:
            try:
                tmp = tempfile.NamedTemporaryFile(delete=False, suffix=".docx")
                tmp.write(raw)
                tmp.flush()
                tmp.close()
                
                d = docx.Document(tmp.name)
                paragraphs = [p.text for p in d.paragraphs if p.text.strip()]
                content = "\n".join(paragraphs)
                
                file_info["paragraphs"] = len(paragraphs)
                
                try:
                    os.unlink(tmp.name)
                except Exception:
                    pass
                    
            except Exception as e:
                content = f"[خطأ في قراءة DOCX: {e}]"
                file_info["error"] = str(e)
    
    elif ext in [".txt", ".md"]:
        try:
            content = raw.decode("utf-8", errors="ignore")
        except Exception as e:
            content = f"[خطأ في قراءة النص: {e}]"
    
    elif ext in [".xls", ".xlsx", ".csv"]:
        try:
            if ext == ".csv":
                df = pd.read_csv(io.BytesIO(raw))
            else:
                if openpyxl is None:
                    content = "[قراءة Excel غير متاحة: تثبيت openpyxl مطلوب]"
                else:
                    df = pd.read_excel(io.BytesIO(raw))
            
            if isinstance(df, pd.DataFrame):
                content = df.fillna("").to_string()
                file_info["rows"] = len(df)
                file_info["columns"] = len(df.columns)
                
        except Exception as e:
            content = f"[خطأ في قراءة الجدول: {e}]"
            file_info["error"] = str(e)
    
    elif ext in [".jpg", ".jpeg", ".png", ".bmp", ".tiff", ".gif"]:
        content = extract_text_from_image(raw)
        file_info["type"] = "image"
    
    elif ext == ".pptx":
        if Presentation is None:
            content = "[قراءة PowerPoint غير متاحة: تثبيت python-pptx مطلوب]"
        else:
            try:
                tmp = tempfile.NamedTemporaryFile(delete=False, suffix=".pptx")
                tmp.write(raw)
                tmp.flush()
                tmp.close()
                
                prs = Presentation(tmp.name)
                slides_text = []
                
                for i, slide in enumerate(prs.slides):
                    slide_text = f"--- شريحة {i+1} ---\n"
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text.strip():
                            slide_text += shape.text + "\n"
                    if slide_text.strip() != f"--- شريحة {i+1} ---":
                        slides_text.append(slide_text)
                
                content = "\n\n".join(slides_text)
                file_info["slides"] = len(prs.slides)
                file_info["extracted_slides"] = len(slides_text)
                
                try:
                    os.unlink(tmp.name)
                except Exception:
                    pass
                    
            except Exception as e:
                content = f"[خطأ في قراءة PowerPoint: {e}]"
                file_info["error"] = str(e)
    
    else:
        try:
            # محاولة قراءة كنص
            content = raw.decode("utf-8", errors="ignore")
        except Exception:
            content = "[نوع ملف غير مدعوم أو محتوى ثنائي]"
    
    return content, file_info

def fallback_clean_text(text: str) -> str:
    """تنظيف محسن للنصوص العربية"""
    if not text:
        return ""
    
    # إزالة التشكيل الزائد (اختياري)
    text = re.sub(r'[\u064B-\u0652]', '', text)
    
    # توحيد المسافات والفواصل
    text = re.sub(r"\s+", " ", text)
    text = re.sub(r"[\r\n]+", "\n", text)
    
    # إزالة الأحرف الخاصة الزائدة مع الاحتفاظ بعلامات الترقيم المهمة
    text = re.sub(r'[^\w\s\u0600-\u06FF\u0750-\u077F\u08A0-\u08FF\uFE70-\uFEFF\u0750-\u077F،؛؟!.()،]', ' ', text)
    
    # تنظيف المسافات الإضافية
    text = re.sub(r'\s+', ' ', text)
    
    return text.strip()

def fallback_chunk_text(text: str, chunk_size: int = 300, overlap: int = 50) -> List[str]:
    """تقسيم النص مع مراعاة اللغة العربية والسياق"""
    if not text:
        return []
    
    # تقسيم بالجمل أولاً
    sentences = fallback_sentences_from_text(text)
    
    if len(sentences) <= 3:
        return [text]
    
    chunks = []
    current_chunk = []
    current_length = 0
    
    for sentence in sentences:
        sentence_words = len(sentence.split())
        
        # إذا كانت الجملة طويلة جداً، قسمها
        if sentence_words > chunk_size:
            if current_chunk:
                chunks.append(" ".join(current_chunk))
                current_chunk = []
                current_length = 0
            
            # تقسيم الجملة الطويلة
            words = sentence.split()
            for i in range(0, len(words), chunk_size - overlap):
                chunk_words = words[i:i + chunk_size]
                chunks.append(" ".join(chunk_words))
        
        # إذا كانت إضافة الجملة ستتجاوز الحد الأقصى
        elif current_length + sentence_words > chunk_size and current_chunk:
            chunks.append(" ".join(current_chunk))
            # الاحتفاظ ببعض الجمل للتداخل
            overlap_sentences = current_chunk[-min(2, len(current_chunk)):]
            current_chunk = overlap_sentences + [sentence]
            current_length = sum(len(s.split()) for s in current_chunk)
        else:
            current_chunk.append(sentence)
            current_length += sentence_words
    
    # إضافة القطعة الأخيرة
    if current_chunk:
        chunks.append(" ".join(current_chunk))
    
    return [chunk for chunk in chunks if len(chunk.strip()) > 20]

from typing import List, Dict, Any
import streamlit as st
import re
from collections import Counter
import math

def fallback_sentences_from_text(text: str) -> List[str]:
    """تقسيم النص لجمل مع مراعاة اللغة العربية"""
    if not text:
        return []
    
    try:
        import nltk
        nltk.download('punkt', quiet=True)
        from nltk.tokenize import sent_tokenize
        sentences = sent_tokenize(text)
        return [s.strip() for s in sentences if len(s.strip()) > 15]
    except Exception:
        # تقسيم احتياطي محسن للعربية
        # إضافة نقاط التقسيم العربية
        text = text.replace('؟', '؟\n')
        text = text.replace('!', '!\n') 
        text = text.replace('.', '.\n')
        text = text.replace('؛', '؛\n')
        text = text.replace(':', ':\n')  # إضافة النقطتين
        text = text.replace('،', '،\n')  # الفاصلة العربية للجمل الطويلة
        
        sentences = [s.strip() for s in text.split('\n') if s.strip()]
        return [s for s in sentences if len(s.strip()) > 15]

def clean_arabic_text(text: str) -> str:
    """تنظيف النص العربي"""
    if not text:
        return ""
    
    # إزالة التشكيل
    text = re.sub(r'[ًٌٍَُِّْ]', '', text)
    
    # توحيد الأحرف المتشابهة
    text = text.replace('أ', 'ا').replace('إ', 'ا').replace('آ', 'ا')
    text = text.replace('ة', 'ه')
    text = text.replace('ى', 'ي')
    
    # إزالة الأرقام والرموز غير المرغوبة
    text = re.sub(r'[0-9]+', '', text)
    text = re.sub(r'[^\u0600-\u06FF\s]', ' ', text)
    
    # إزالة المسافات الزائدة
    text = ' '.join(text.split())
    
    return text.strip()

def fallback_tfidf_sentence_ranking(document_texts: List[str], top_k_sentences_per_doc: int = 3) -> Dict[str, List[Dict]]:
    """ترتيب الجمل حسب الأهمية مع تحسينات للعربية"""
    if not document_texts:
        return {}
    
    try:
        # جمع كل الجمل من كل الوثائق
        all_sentences = []
        doc_sentence_map = {}
        
        for doc_idx, doc_text in enumerate(document_texts):
            sentences = fallback_sentences_from_text(doc_text)
            doc_id = f"doc_{doc_idx + 1}"
            doc_sentence_map[doc_id] = sentences
            all_sentences.extend(sentences)
        
        if not all_sentences:
            return {}
        
        # حساب TF-IDF بسيط
        # 1. حساب تكرار الكلمات
        word_doc_count = Counter()
        sentence_word_counts = []
        
        for sentence in all_sentences:
            clean_sentence = clean_arabic_text(sentence)
            words = clean_sentence.split()
            word_count = Counter(words)
            sentence_word_counts.append(word_count)
            
            # حساب عدد الوثائق التي تحتوي على كل كلمة
            for word in set(words):
                if len(word) > 2:  # تجاهل الكلمات القصيرة
                    word_doc_count[word] += 1
        
        # 2. حساب نقاط TF-IDF للجمل
        sentence_scores = []
        total_docs = len(all_sentences)
        
        for word_count in sentence_word_counts:
            score = 0
            total_words = sum(word_count.values())
            
            if total_words > 0:
                for word, count in word_count.items():
                    if len(word) > 2 and word_doc_count[word] > 0:
                        tf = count / total_words
                        idf = math.log(total_docs / word_doc_count[word])
                        score += tf * idf
            
            sentence_scores.append(score)
        
        # 3. ترتيب الجمل لكل وثيقة
        results = {}
        sentence_idx = 0
        
        for doc_id, sentences in doc_sentence_map.items():
            doc_sentences_with_scores = []
            
            for sentence in sentences:
                if sentence_idx < len(sentence_scores):
                    doc_sentences_with_scores.append({
                        'text': sentence,
                        'score': sentence_scores[sentence_idx],
                        'length': len(sentence)
                    })
                sentence_idx += 1
            
            # ترتيب حسب النقاط
            doc_sentences_with_scores.sort(key=lambda x: x['score'], reverse=True)
            
            # أخذ أفضل الجمل
            results[doc_id] = doc_sentences_with_scores[:top_k_sentences_per_doc]
        
        return results
    
    except Exception as e:
        st.error(f"خطأ في ترتيب الجمل: {e}")
        return {}

def process_documents(documents_list=None):
    """معالجة قائمة الوثائق وتنظيفها"""
    try:
        # استخدام المتغير المُمرر أو البحث عن document_texts
        if documents_list is None:
            if 'document_texts' in st.session_state:
                documents_list = st.session_state.document_texts
            elif 'document_texts' in globals():
                documents_list = document_texts
            else:
                st.error("لم يتم العثور على أي وثائق للمعالجة.")
                return []
        
        if not documents_list:
            st.warning("قائمة الوثائق فارغة.")
            return []
            
        processed_docs = []
        progress_bar = st.progress(0)
        
        for i, doc in enumerate(documents_list):
            # تنظيف النص
            clean_doc = doc.strip().replace("\n", " ").replace("\r", "")
            clean_doc = clean_arabic_text(clean_doc)  # تنظيف النص العربي
            doc_id = f"doc_{i+1}"
            
            # تقسيم إلى جمل
            sentences = fallback_sentences_from_text(doc)
            
            processed_docs.append({
                "id": doc_id,
                "text": clean_doc,
                "original_text": doc,
                "sentences": sentences,
                "length": len(clean_doc),
                "sentence_count": len(sentences)
            })
            
            # تحديث شريط التقدم
            progress_bar.progress((i + 1) / len(documents_list))
            st.write(f"✅ تم معالجة الوثيقة {doc_id} - الطول: {len(clean_doc)} حرف - الجمل: {len(sentences)}")
            
        st.success(f"🎉 تم معالجة {len(processed_docs)} وثيقة بنجاح!")
        
        # حفظ في session state
        st.session_state.processed_docs = processed_docs
        
        return processed_docs
        
    except Exception as e:
        st.error(f"❌ خطأ في معالجة الوثائق: {e}")
        return []

def extract_key_sentences(processed_docs: List[Dict], top_k: int = 3) -> Dict:
    """استخراج أهم الجمل من كل وثيقة"""
    if not processed_docs:
        return {}
    
    # تحضير قائمة النصوص للتحليل
    document_texts = [doc['original_text'] for doc in processed_docs]
    
    # تحليل وترتيب الجمل
    ranked_sentences = fallback_tfidf_sentence_ranking(document_texts, top_k)
    
    return ranked_sentences

def create_searchable_index(processed_docs: List[Dict]) -> Dict:
    """إنشاء فهرس للبحث"""
    search_index = {
        'documents': {},
        'words': {},
        'sentences': []
    }
    
    for doc in processed_docs:
        doc_id = doc['id']
        clean_text = doc['text']
        sentences = doc['sentences']
        
        # فهرسة الوثيقة
        search_index['documents'][doc_id] = {
            'text': clean_text,
            'sentences': sentences,
            'word_count': len(clean_text.split())
        }
        
        # فهرسة الكلمات
        words = clean_text.split()
        for word in words:
            if len(word) > 2:
                if word not in search_index['words']:
                    search_index['words'][word] = []
                search_index['words'][word].append(doc_id)
        
        # فهرسة الجمل
        for sentence in sentences:
            clean_sentence = clean_arabic_text(sentence)
            search_index['sentences'].append({
                'text': sentence,
                'clean_text': clean_sentence,
                'doc_id': doc_id,
                'words': clean_sentence.split()
            })
    
    return search_index

# مثال للاستخدام
def main():
    st.title("🌍 منصة RAG المتطورة للوثائق العربية")
    
    # تحميل الوثائق
    if 'document_texts' not in st.session_state:
        st.session_state.document_texts = [
            "هذا نص تجريبي باللغة العربية. يحتوي على عدة جمل للاختبار.",
            "النص الثاني يتضمن معلومات مختلفة. هدفه اختبار النظام.",
            "الوثيقة الثالثة تحتوي على محتوى متنوع ومفيد للغاية."
        ]
    
    if st.button("معالجة الوثائق"):
        processed_docs = process_documents()
        
        if processed_docs:
            # استخراج الجمل المهمة
            key_sentences = extract_key_sentences(processed_docs)
            
            st.subheader("🔍 أهم الجمل من كل وثيقة:")
            for doc_id, sentences in key_sentences.items():
                st.write(f"**{doc_id}:**")
                for i, sent_info in enumerate(sentences, 1):
                    st.write(f"  {i}. {sent_info['text']} (نقاط: {sent_info['score']:.3f})")
                st.write("")
            
            # إنشاء فهرس البحث
            search_index = create_searchable_index(processed_docs)
            st.session_state.search_index = search_index
            
            st.success("تم إنشاء فهرس البحث بنجاح!")

if __name__ == "__main__":
    main()
